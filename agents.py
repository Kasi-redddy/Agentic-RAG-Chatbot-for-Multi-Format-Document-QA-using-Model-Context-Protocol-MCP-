import PyPDF2
import pandas as pd
from docx import Document
from pptx import Presentation
from sentence_transformers import SentenceTransformer
import faiss
import numpy as np
import ollama
from pydantic import BaseModel
from typing import Dict, Any, List, Optional
import uuid
import os
import logging
import nltk
from nltk.tokenize import sent_tokenize

# Set up logging for debugging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# Download NLTK punkt tokenizer if not present
try:
    nltk.data.find("tokenizers/punkt")
except LookupError:
    nltk.download("punkt")

# MCP Message Class matching the specified format
class MCPMessage(BaseModel):
    sender: str
    receiver: str
    type: str # e.g., "DOCUMENT_UPLOAD", "DOCUMENT_PARSED", "CONTEXT_REQUEST", "CONTEXT_RESPONSE", "ANSWER_RESPONSE", "ERROR"
    trace_id: str
    payload: Dict[str, Any]

# Document Parser Utility
def parse_document(filename: str, content: bytes) -> str:
    """Parse documents (PDF, CSV, DOCX, PPTX, TXT, Markdown) from bytes."""
    try:
        ext = os.path.splitext(filename)[1].lower()
        logger.info(f"Parsing document: {filename} with extension {ext}")

        if ext == ".pdf":
            from io import BytesIO
            reader = PyPDF2.PdfReader(BytesIO(content))
            text = " ".join(page.extract_text() or "" for page in reader.pages)
            if not text.strip():
                raise ValueError("No text extracted from PDF. It might be an image-based PDF or empty.")
            return text
        
        elif ext == ".csv":
            from io import StringIO
            for encoding in ["utf-8", "latin-1", "iso-8859-1"]:
                try:
                    df = pd.read_csv(StringIO(content.decode(encoding)), on_bad_lines="skip")
                    df = df.astype(str)
                    text = f"Columns: {', '.join(df.columns)}\n"
                    for _, row in df.iterrows():
                        text += "; ".join(f"{col}: {val}" for col, val in row.items()) + "\n"
                    return text
                except UnicodeDecodeError:
                    continue
            raise ValueError("Failed to decode CSV with supported encodings.")
        
        elif ext == ".docx":
            from io import BytesIO
            doc = Document(BytesIO(content))
            text = " ".join(paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip())
            if not text:
                raise ValueError("No text extracted from DOCX. The document might be empty or corrupt.")
            return text
        
        elif ext == ".pptx":
            from io import BytesIO
            prs = Presentation(BytesIO(content))
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.text.strip():
                                    text_runs.append(run.text)
            text = " ".join(text_runs)
            if not text.strip():
                raise ValueError("No text extracted from PPTX. The presentation might be empty or contain only images.")
            return text

        elif ext in [".txt", ".md"]:
            text = content.decode("utf-8", errors="ignore")
            if not text.strip():
                raise ValueError("No text extracted from TXT/Markdown. The file might be empty.")
            return text
        
        raise ValueError(f"Unsupported file format: {ext}. Please upload PDF, CSV, DOCX, PPTX, TXT, or Markdown files.")
    
    except Exception as e:
        logger.error(f"Error parsing {filename}: {str(e)}", exc_info=True)
        raise

# Represents a chunk of text with its source information
class DocumentChunk(BaseModel):
    doc_id: str
    filename: str
    chunk_text: str
    chunk_index: int # To help identify position within the document

# Vector Store for embeddings and retrieval
class VectorStore:
    def __init__(self):
        self.model = SentenceTransformer("all-MiniLM-L6-v2") 
        self.index = faiss.IndexFlatL2(self.model.get_sentence_embedding_dimension())
        self.doc_chunks: List[DocumentChunk] = [] 
        self.embedding_cache = {} 
        logger.info(f"VectorStore initialized with embedding dimension: {self.model.get_sentence_embedding_dimension()}")

    def add_document(self, doc_id: str, filename: str, content: str):
        """Add document with sentence-based chunking and store metadata."""
        try:
            sentences = sent_tokenize(content)
            chunks = []
            current_chunk = ""
            chunk_idx = 0
            
            for sentence in sentences:
                if len(current_chunk) + len(sentence) + (1 if current_chunk else 0) < 500: 
                    current_chunk += (" " if current_chunk else "") + sentence
                else:
                    if current_chunk:
                        chunks.append(DocumentChunk(doc_id=doc_id, filename=filename, chunk_text=current_chunk.strip(), chunk_index=chunk_idx))
                        chunk_idx += 1
                    current_chunk = sentence
            if current_chunk:
                chunks.append(DocumentChunk(doc_id=doc_id, filename=filename, chunk_text=current_chunk.strip(), chunk_index=chunk_idx))

            embeddings = []
            for chunk_obj in chunks:
                if chunk_obj.chunk_text in self.embedding_cache:
                    embeddings.append(self.embedding_cache[chunk_obj.chunk_text])
                else:
                    embedding = self.model.encode([chunk_obj.chunk_text], show_progress_bar=False)[0]
                    embedding = embedding / np.linalg.norm(embedding) 
                    self.embedding_cache[chunk_obj.chunk_text] = embedding
                    embeddings.append(embedding)
            
            if embeddings:
                self.index.add(np.array(embeddings, dtype=np.float32))
                self.doc_chunks.extend(chunks)
                logger.info(f"Added document {filename} (ID: {doc_id}) with {len(chunks)} chunks to VectorStore.")
            else:
                logger.warning(f"No chunks generated for document {filename} (ID: {doc_id}). Content might be too short.")
        except Exception as e:
            logger.error(f"Error adding document {filename} (ID: {doc_id}): {str(e)}", exc_info=True)
            raise

    def search(self, query: str, top_k: int = 5) -> List[DocumentChunk]:
        """Search for top_k relevant chunks."""
        try:
            if self.index.ntotal == 0:
                logger.warning("Vector store is empty. No documents to search.")
                return []

            query_embedding = self.model.encode([query])[0]
            query_embedding = query_embedding / np.linalg.norm(query_embedding)

            distances, indices = self.index.search(np.array([query_embedding], dtype=np.float32), top_k)
            
            results = []
            for i in indices[0]:
                if 0 <= i < len(self.doc_chunks):
                    results.append(self.doc_chunks[i])
            
            logger.info(f"Retrieved {len(results)} chunks for query: '{query}'")
            return results
        except Exception as e:
            logger.error(f"Error searching query '{query}': {str(e)}", exc_info=True)
            raise

# Ingestion Agent
class IngestionAgent:
    def __init__(self):
        self.parsed_docs = {} 
        logger.info("IngestionAgent initialized")

    def process_document(self, filename: str, content: bytes, trace_id: str) -> MCPMessage:
        """Parse document and return MCP message."""
        try:
            content_text = parse_document(filename, content)
            doc_id = str(uuid.uuid4())
            self.parsed_docs[doc_id] = {"filename": filename, "content": content_text}
            logger.info(f"IngestionAgent processed document {filename} with ID {doc_id}")
            return MCPMessage(
                sender="IngestionAgent",
                receiver="RetrievalAgent",
                type="DOCUMENT_PARSED",
                trace_id=trace_id,
                payload={"doc_id": doc_id, "content": content_text, "filename": filename}
            )
        except Exception as e:
            logger.error(f"IngestionAgent error for {filename}: {str(e)}", exc_info=True)
            return MCPMessage(
                sender="IngestionAgent",
                receiver="CoordinatorAgent",
                type="ERROR",
                trace_id=trace_id,
                payload={"error": f"Failed to parse document {filename}: {str(e)}"}
            )

# Retrieval Agent
class RetrievalAgent:
    def __init__(self):
        self.vector_store = VectorStore()
        logger.info("RetrievalAgent initialized")

    def process_document_for_indexing(self, message: MCPMessage) -> MCPMessage:
        """Index document content from an IngestionAgent message in vector store."""
        try:
            doc_id = message.payload["doc_id"]
            content = message.payload["content"]
            filename = message.payload["filename"]
            self.vector_store.add_document(doc_id, filename, content)
            logger.info(f"RetrievalAgent indexed document {filename} (ID: {doc_id})")
            return MCPMessage(
                sender="RetrievalAgent",
                receiver="CoordinatorAgent",
                type="DOCUMENT_INDEXED",
                trace_id=message.trace_id,
                payload={"doc_id": doc_id, "filename": filename}
            )
        except Exception as e:
            logger.error(f"RetrievalAgent error during indexing: {str(e)}", exc_info=True)
            return MCPMessage(
                sender="RetrievalAgent",
                receiver="CoordinatorAgent",
                type="ERROR",
                trace_id=message.trace_id,
                payload={"error": f"Failed to index document: {str(e)}"}
            )

    def retrieve_context(self, message: MCPMessage) -> MCPMessage:
        """Retrieve relevant chunks for a query from a CoordinatorAgent message."""
        try:
            query = message.payload["query"]
            top_chunks = self.vector_store.search(query)
            
            formatted_context_for_mcp = []
            for chunk in top_chunks:
                formatted_context_for_mcp.append({
                    "text": chunk.chunk_text,
                    "source_info": f"{chunk.filename} (chunk {chunk.chunk_index + 1})"
                })
            
            logger.info(f"RetrievalAgent retrieved {len(formatted_context_for_mcp)} chunks for query: '{query}'")
            
            return MCPMessage(
                sender="RetrievalAgent",
                receiver="LLMResponseAgent",
                type="CONTEXT_RESPONSE",
                trace_id=message.trace_id,
                payload={
                    "retrieved_context": [item["text"] for item in formatted_context_for_mcp],
                    "query": query,
                    "full_context_details": formatted_context_for_mcp
                }
            )
        except Exception as e:
            logger.error(f"RetrievalAgent error retrieving context for query '{message.payload.get('query')}': {str(e)}", exc_info=True)
            return MCPMessage(
                sender="RetrievalAgent",
                receiver="CoordinatorAgent",
                type="ERROR",
                trace_id=message.trace_id,
                payload={"error": f"Failed to retrieve context: {str(e)}"}
            )

# LLM Response Agent
class LLMResponseAgent:
    def __init__(self):
        self.model_name = "llama3"
        logger.info(f"LLMResponseAgent initialized with Ollama model: {self.model_name}")

    def generate_response(self, message: MCPMessage) -> MCPMessage:
        """Generate response using Ollama's LLM based on retrieved context."""
        try:
            query = message.payload["query"]
            context_texts = message.payload["retrieved_context"] 
            full_context_details = message.payload.get("full_context_details", []) 
            
            context = "\n\n".join(context_texts)

            if not context.strip():
                answer = "I could not find relevant information in the uploaded documents to answer your question."
                logger.warning(f"LLMResponseAgent: No context provided for query: {query}")
            else:
                prompt = f"""
                You are a helpful AI assistant. Answer the user's question ONLY based on the provided context.
                If the answer cannot be found in the context, state that you don't have enough information.

                Context:
                {context}

                Question: {query}

                Answer:
                """
                
                logger.info(f"Calling Ollama with prompt for query: {query}")
                response = ollama.generate(
                    model=self.model_name,
                    prompt=prompt,
                    options={"num_predict": 512, "temperature": 0.3}
                )
                answer = response["response"].strip()
                logger.info(f"LLMResponseAgent generated response for query: {query}")

            retrieval_result_mcp_payload = {
                "type": "RETRIEVAL_RESULT",
                "sender": "RetrievalAgent",
                "receiver": "LLMResponseAgent",
                "trace_id": message.trace_id,
                "payload": {
                    "retrieved_context": context_texts,
                    "query": query
                }
            }

            return MCPMessage(
                sender="LLMResponseAgent",
                receiver="CoordinatorAgent",
                type="ANSWER_RESPONSE",
                trace_id=message.trace_id,
                payload={
                    "answer": answer,
                    "retrieval_mcp_output": retrieval_result_mcp_payload,
                    "original_sources_info": full_context_details
                }
            )
        except Exception as e:
            logger.error(f"LLMResponseAgent error generating response for query '{message.payload.get('query')}': {str(e)}", exc_info=True)
            return MCPMessage(
                sender="LLMResponseAgent",
                receiver="CoordinatorAgent",
                type="ERROR",
                trace_id=message.trace_id,
                payload={"error": f"Failed to generate LLM response: {str(e)}"}
            )

# Coordinator Agent
class CoordinatorAgent:
    def __init__(self):
        self.ingestion_agent = IngestionAgent()
        self.retrieval_agent = RetrievalAgent()
        self.llm_response_agent = LLMResponseAgent()
        self.message_log: List[MCPMessage] = []
        logger.info("CoordinatorAgent initialized")

    def _send_message(self, message: MCPMessage):
        """Simulate sending a message by adding it to a log and processing."""
        self.message_log.append(message)
        logger.info(f"MCP Message Sent: Sender={message.sender}, Receiver={message.receiver}, Type={message.type}, TraceID={message.trace_id}")
        
    def process_document_upload(self, filename: str, content: bytes):
        """Initiate document processing workflow."""
        trace_id = f"doc-{uuid.uuid4().hex[:8]}"
        initial_message = MCPMessage(
            sender="UI",
            receiver="CoordinatorAgent",
            type="DOCUMENT_UPLOAD",
            trace_id=trace_id,
            payload={"filename": filename, "content": content}
        )
        self._send_message(initial_message)

        ingestion_msg = self.ingestion_agent.process_document(filename, content, trace_id)
        self._send_message(ingestion_msg)

        if ingestion_msg.type == "DOCUMENT_PARSED":
            retrieval_indexing_msg = self.retrieval_agent.process_document_for_indexing(ingestion_msg)
            self._send_message(retrieval_indexing_msg)
            if retrieval_indexing_msg.type == "ERROR":
                raise Exception(f"Indexing Error: {retrieval_indexing_msg.payload['error']}")
        elif ingestion_msg.type == "ERROR":
            raise Exception(f"Parsing Error: {ingestion_msg.payload['error']}")
        else:
            raise Exception(f"Unexpected message type from IngestionAgent: {ingestion_msg.type}")


    def handle_query(self, query: str) -> tuple[str, Dict[str, Any], List[Dict[str, Any]]]: # Adjusted type hint for original_sources_info
        """Handle user query and orchestrate agents."""
        trace_id = f"query-{uuid.uuid4().hex[:8]}"
        
        query_request_msg = MCPMessage(
            sender="UI",
            receiver="CoordinatorAgent",
            type="QUERY_REQUEST",
            trace_id=trace_id,
            payload={"query": query}
        )
        self._send_message(query_request_msg)

        retrieval_msg = self.retrieval_agent.retrieve_context(query_request_msg)
        self._send_message(retrieval_msg)

        if retrieval_msg.type == "CONTEXT_RESPONSE":
            response_message = self.llm_response_agent.generate_response(retrieval_msg)
            self._send_message(response_message)
            
            if response_message.type == "ANSWER_RESPONSE":
                return (
                    response_message.payload["answer"],
                    response_message.payload["retrieval_mcp_output"],
                    response_message.payload["original_sources_info"]
                )
            elif response_message.type == "ERROR":
                raise Exception(f"LLM Response Error: {response_message.payload['error']}")
            else:
                raise Exception(f"Unexpected message type from LLMResponseAgent: {response_message.type}")
        elif retrieval_msg.type == "ERROR":
            raise Exception(f"Retrieval Error: {retrieval_msg.payload['error']}")
        else:
            raise Exception(f"Unexpected message type from RetrievalAgent: {retrieval_msg.type}")